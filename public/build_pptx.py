#!/usr/bin/env python3
"""Generate SignalSniper AI hackathon presentation (.pptx) for Google Slides import."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BG = RGBColor(0x0A, 0x0A, 0x0A)
SURFACE = RGBColor(0x14, 0x14, 0x14)
WHITE = RGBColor(0xF0, 0xF0, 0xF0)
DIM = RGBColor(0x80, 0x80, 0x80)
MUTED = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x00, 0xE8, 0x7B)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PINK = RGBColor(0xED, 0x64, 0xA6)
BORDER = RGBColor(0x25, 0x25, 0x25)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill=SURFACE, border=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    return sh


def accent_bar(slide, l, t, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.06), h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def txt(slide, l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def multi(slide, l, t, w, h, lines, sz=12, color=DIM):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(sz * 0.6)
    return box


def eyebrow(slide, text):
    txt(slide, Inches(0.9), Inches(0.7), Inches(5), Inches(0.35), text, sz=10, color=GREEN, bold=True, font="Consolas")


def title_block(slide, line1, line2="", line2_color=WHITE):
    txt(slide, Inches(0.9), Inches(1.2), Inches(10), Inches(0.9), line1, sz=42, color=WHITE, bold=True, font="Georgia")
    if line2:
        txt(slide, Inches(0.9), Inches(2.0), Inches(10), Inches(0.8), line2, sz=38, color=line2_color, bold=True, font="Georgia")


def sub(slide, text, w=Inches(6)):
    txt(slide, Inches(0.9), Inches(2.8), w, Inches(0.9), text, sz=15, color=DIM)


# ═══════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
txt(s, Inches(0.9), Inches(1.0), Inches(5), Inches(0.5), "◎  SignalSniper AI", sz=20, color=WHITE, bold=True)
txt(s, Inches(0.9), Inches(2.0), Inches(5), Inches(0.35), "AMD HACKATHON 2026", sz=10, color=GREEN, bold=True, font="Consolas")
txt(s, Inches(0.9), Inches(2.8), Inches(10), Inches(1), "Forex signals,", sz=54, color=WHITE, bold=True, font="Georgia")
txt(s, Inches(0.9), Inches(3.8), Inches(10), Inches(1), "explained by AI.", sz=54, color=GREEN, bold=True, font="Georgia")
txt(s, Inches(0.9), Inches(5.0), Inches(7), Inches(0.9),
    "A dual-AI pipeline that transforms raw trading signals into structured, trader-grade analysis — powered by Fireworks AI and AMD GPU-accelerated Gemma inference.",
    sz=15, color=DIM)
rect(s, Inches(0.9), Inches(6.2), Inches(2), Inches(0.35), SURFACE, BORDER)
txt(s, Inches(1.1), Inches(6.22), Inches(1.8), Inches(0.3), "● Live Demo Ready", sz=9, color=GREEN, font="Consolas")
rect(s, Inches(3.1), Inches(6.2), Inches(3.6), Inches(0.35), SURFACE, BORDER)
txt(s, Inches(3.3), Inches(6.22), Inches(3.4), Inches(0.3), "Next.js · Fireworks AI · Gemma · vLLM", sz=9, color=MUTED, font="Consolas")

# ═══════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
eyebrow(s, "THE PROBLEM")
title_block(s, "Traders drown in signals", "but starve for understanding.", DIM)
sub(s, "")

cards_data = [
    ("73%", "OF RETAIL TRADERS LOSE MONEY", ROSE, "Most follow signals blindly without understanding the reasoning — leading to poor execution and emotional exits."),
    ("BUY", "THAT'S THE ENTIRE SIGNAL", AMBER, "Traditional signal services output a single word. No entry context, no risk assessment, no market sentiment."),
    ("0s", "EXPLAINABILITY", PURPLE, "Traders have no way to understand why a signal was generated. Technical analysis is reduced to a badge."),
]
for i, (num, label, nc, desc) in enumerate(cards_data):
    left = Inches(0.9 + i * 3.9)
    rect(s, left, Inches(3.4), Inches(3.5), Inches(3.2), SURFACE, BORDER)
    txt(s, left + Inches(0.3), Inches(3.7), Inches(2.9), Inches(0.9), num, sz=42, color=nc, bold=True, font="Georgia")
    txt(s, left + Inches(0.3), Inches(4.5), Inches(2.9), Inches(0.3), label, sz=8, color=MUTED, bold=True, font="Consolas")
    txt(s, left + Inches(0.3), Inches(5.0), Inches(2.9), Inches(1.2), desc, sz=11, color=DIM)

# ═══════════════════════════════════════
# SLIDE 3: THE SOLUTION
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
eyebrow(s, "THE SOLUTION")
title_block(s, "One signal.", "Full trade thesis.", GREEN)

features = [
    ("📊", "14 Technical Indicators", "RSI, MACD, Bollinger Bands, EMA, SMA, Stochastic, ATR, ADX — all fed into the AI prompt context."),
    ("🎯", "Trader-Grade Output", "Entry price, stop loss, take profit, R:R ratio, support/resistance levels, key pivot points."),
    ("🧠", "AI-Explained Reasoning", "Market sentiment, numbered reasons, risk warnings — not just BUY but why you should."),
]
for i, (icon, title, desc) in enumerate(features):
    top = Inches(3.3 + i * 1.15)
    rect(s, Inches(0.9), top, Inches(5.8), Inches(0.95), SURFACE, BORDER)
    txt(s, Inches(1.15), top + Inches(0.1), Inches(0.4), Inches(0.4), icon, sz=16, align=PP_ALIGN.CENTER)
    txt(s, Inches(1.6), top + Inches(0.1), Inches(4.8), Inches(0.35), title, sz=14, color=WHITE, bold=True)
    txt(s, Inches(1.6), top + Inches(0.45), Inches(4.8), Inches(0.4), desc, sz=11, color=DIM)

# Code block
rect(s, Inches(7.2), Inches(3.3), Inches(5.2), Inches(3.6), SURFACE, BORDER)
code = [
    '// Structured AI output — not a black box',
    '{',
    '  "recommendation": "SELL",',
    '  "confidence": 68,',
    '  "trend": "BEARISH",',
    '  "risk": "MEDIUM",',
    '  "entryPrice": 1.1035,',
    '  "stopLoss": 1.1150,',
    '  "takeProfit": 1.0900,',
    '  "riskRewardRatio": "1:2",',
    '  "marketSentiment": "...",',
    '  "reasons": ["..."],',
    '  "warnings": ["..."]',
    '}',
]
multi(s, Inches(7.5), Inches(3.5), Inches(4.8), Inches(3.2), code, sz=10, color=DIM)

# ═══════════════════════════════════════
# SLIDE 4: DUAL AI PIPELINE
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
eyebrow(s, "ARCHITECTURE")
title_block(s, "Dual AI Pipeline")
sub(s, "Two specialized AI layers — one for structured analysis, one for narrative explanation.", Inches(6))

# Pipeline 1
rect(s, Inches(0.9), Inches(3.5), Inches(5.5), Inches(3.5), SURFACE, CYAN)
accent_bar(s, Inches(0.9), Inches(3.5), Inches(3.5), CYAN)
rect(s, Inches(1.2), Inches(3.7), Inches(2.5), Inches(0.3), SURFACE, CYAN)
txt(s, Inches(1.4), Inches(3.72), Inches(2.3), Inches(0.25), "● Pipeline 1 — Signal Generation", sz=9, color=CYAN, font="Consolas")
txt(s, Inches(1.2), Inches(4.1), Inches(5.0), Inches(0.4), "Fireworks AI → Structured JSON", sz=16, color=CYAN, bold=True)
steps1 = [
    "Live Rates API  →  14 Indicators",
    "Prompt Builder  →  Gemma 3N 27B",
    "JSON Extractor  →  Cache (1hr TTL)",
]
multi(s, Inches(1.2), Inches(4.6), Inches(5.0), Inches(1.2), steps1, sz=11, color=DIM)
txt(s, Inches(1.2), Inches(5.8), Inches(5.0), Inches(0.8),
    "Fallback chain: gemma3n-27b → gpt-oss-120b.\nIf primary fails, secondary takes over automatically.",
    sz=10, color=MUTED)

# Pipeline 2
rect(s, Inches(6.9), Inches(3.5), Inches(5.5), Inches(3.5), SURFACE, PURPLE)
accent_bar(s, Inches(6.9), Inches(3.5), Inches(3.5), PURPLE)
rect(s, Inches(7.2), Inches(3.7), Inches(2.5), Inches(0.3), SURFACE, PURPLE)
txt(s, Inches(7.4), Inches(3.72), Inches(2.3), Inches(0.25), "● Pipeline 2 — Deep Explanation", sz=9, color=PURPLE, font="Consolas")
txt(s, Inches(7.2), Inches(4.1), Inches(5.0), Inches(0.4), "AMD GPU + vLLM → Narrative", sz=16, color=PURPLE, bold=True)
steps2 = [
    "Signal Data  →  Prompt Template",
    "vLLM Server  →  AMD GPU (ROCm)",
    "Gemma Model  →  Analysis Text",
]
multi(s, Inches(7.2), Inches(4.6), Inches(5.0), Inches(1.2), steps2, sz=11, color=DIM)
txt(s, Inches(7.2), Inches(5.8), Inches(5.0), Inches(0.8),
    "Mock ↔ Real toggle: one env var switches\nbetween mock and live AMD GPU inference.\nZero code changes.",
    sz=10, color=MUTED)

# ═══════════════════════════════════════
# SLIDE 5: DATA FLOW
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
eyebrow(s, "DATA FLOW")
title_block(s, "End-to-End Request Path")
sub(s, "From browser click to AI inference and back — a complete signal analysis in under 3 seconds.", Inches(7))

steps = [
    ("1", 'User clicks "Analyze" on a currency pair', "POST /api/analyze with pair name + live exchange rate", GREEN),
    ("2", "Server checks in-memory cache (1hr TTL)", "Fresh data returned instantly; otherwise proceeds to Fireworks AI", CYAN),
    ("3", "Fireworks AI receives prompt with 14 indicators", "RSI, MACD, Bollinger, EMA, SMA, Stochastic, ATR, ADX, volume", PURPLE),
    ("4", "Gemma AI deep-dive via AMD GPU (vLLM)", "Signal data → prompt template → self-hosted Gemma → narrative", AMBER),
    ("5", "Dashboard renders full trade card", "Confidence, trend/risk, entry/SL/TP, key levels, sentiment, reasons", BLUE),
]
for i, (num, title, desc, nc) in enumerate(steps):
    top = Inches(3.5 + i * 0.75)
    # Number circle
    rect(s, Inches(1.0), top, Inches(0.45), Inches(0.45), SURFACE, nc)
    txt(s, Inches(1.0), top + Inches(0.05), Inches(0.45), Inches(0.35), num, sz=16, color=nc, bold=True, align=PP_ALIGN.CENTER)
    # Title + desc
    txt(s, Inches(1.7), top + Inches(0.02), Inches(5), Inches(0.35), title, sz=13, color=WHITE, bold=True)
    txt(s, Inches(1.7), top + Inches(0.35), Inches(7), Inches(0.3), desc, sz=10, color=DIM)

# ═══════════════════════════════════════
# SLIDE 6: TECH STACK
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
eyebrow(s, "TECH STACK")
title_block(s, "Built With")

# Frontend
txt(s, Inches(0.9), Inches(3.0), Inches(3), Inches(0.3), "FRONTEND", sz=9, color=MUTED, bold=True, font="Consolas")
fe_tags = ["Next.js 16", "React 19", "TypeScript", "Tailwind CSS 4", "Canvas API"]
for i, tag in enumerate(fe_tags):
    x = Inches(0.9 + i * 1.7)
    rect(s, x, Inches(3.3), Inches(1.5), Inches(0.35), SURFACE, BORDER)
    txt(s, x + Inches(0.12), Inches(3.33), Inches(1.3), Inches(0.3), tag, sz=9, color=WHITE, font="Consolas")

# AI / Inference
txt(s, Inches(0.9), Inches(4.0), Inches(3), Inches(0.3), "AI / INFERENCE", sz=9, color=MUTED, bold=True, font="Consolas")
ai_tags = ["Fireworks AI", "Gemma 3N 27B", "GPT-OSS 120B", "vLLM Server", "ROCm / AMD GPU"]
ai_colors = [CYAN, PURPLE, PURPLE, PINK, PINK]
for i, (tag, tc) in enumerate(zip(ai_tags, ai_colors)):
    x = Inches(0.9 + i * 2.1)
    rect(s, x, Inches(4.3), Inches(1.9), Inches(0.35), SURFACE, tc)
    txt(s, x + Inches(0.12), Inches(4.33), Inches(1.7), Inches(0.3), tag, sz=9, color=WHITE, font="Consolas")

# Data
txt(s, Inches(0.9), Inches(5.0), Inches(3), Inches(0.3), "DATA", sz=9, color=MUTED, bold=True, font="Consolas")
data_tags = ["ExchangeRate API", "In-Memory Cache", "Promise.all Parallel"]
for i, tag in enumerate(data_tags):
    x = Inches(0.9 + i * 2.4)
    rect(s, x, Inches(5.3), Inches(2.2), Inches(0.35), SURFACE, BORDER)
    txt(s, x + Inches(0.12), Inches(5.33), Inches(2.0), Inches(0.3), tag, sz=9, color=WHITE, font="Consolas")

# Stats
rect(s, Inches(9.5), Inches(3.0), Inches(3.0), Inches(3.5), SURFACE, BORDER)
txt(s, Inches(9.5), Inches(3.2), Inches(3.0), Inches(0.8), "16", sz=48, color=GREEN, bold=True, font="Georgia", align=PP_ALIGN.CENTER)
txt(s, Inches(9.5), Inches(3.9), Inches(3.0), Inches(0.3), "SOURCE FILES", sz=9, color=MUTED, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
txt(s, Inches(9.5), Inches(4.5), Inches(3.0), Inches(0.8), "4", sz=48, color=CYAN, bold=True, font="Georgia", align=PP_ALIGN.CENTER)
txt(s, Inches(9.5), Inches(5.2), Inches(3.0), Inches(0.3), "API ENDPOINTS", sz=9, color=MUTED, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
txt(s, Inches(9.5), Inches(5.6), Inches(3.0), Inches(0.8), "3", sz=48, color=PURPLE, bold=True, font="Georgia", align=PP_ALIGN.CENTER)
txt(s, Inches(9.5), Inches(6.2), Inches(3.0), Inches(0.3), "FOREX PAIRS", sz=9, color=MUTED, bold=True, font="Consolas", align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 7: AMD GPU INTEGRATION
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
eyebrow(s, "AMD INTEGRATION")
title_block(s, "GPU-Accelerated Inference", "on AMD Hardware", DIM)

# vLLM + ROCm
rect(s, Inches(0.9), Inches(3.3), Inches(5.5), Inches(2.0), SURFACE, BORDER)
rect(s, Inches(1.2), Inches(3.5), Inches(1.8), Inches(0.3), SURFACE, PINK)
txt(s, Inches(1.4), Inches(3.52), Inches(1.6), Inches(0.25), "● vLLM + ROCm", sz=9, color=PINK, font="Consolas")
txt(s, Inches(1.2), Inches(3.95), Inches(5.0), Inches(1.2),
    "The Gemma deep-explanation pipeline runs on a self-hosted vLLM server with AMD ROCm acceleration. The OpenAI-compatible /v1/chat/completions endpoint serves the Gemma model for narrative analysis.",
    sz=12, color=DIM)

# Zero-Config Toggle
rect(s, Inches(0.9), Inches(5.5), Inches(5.5), Inches(1.7), SURFACE, BORDER)
rect(s, Inches(1.2), Inches(5.7), Inches(2.1), Inches(0.3), SURFACE, GREEN)
txt(s, Inches(1.4), Inches(5.72), Inches(1.9), Inches(0.25), "● Zero-Config Toggle", sz=9, color=GREEN, font="Consolas")
txt(s, Inches(1.2), Inches(6.15), Inches(5.0), Inches(0.9),
    "A single environment variable switches between mock and real AMD GPU inference. The API route, frontend, and UI remain identical — no code changes needed.",
    sz=12, color=DIM)

# Why AMD
rect(s, Inches(6.9), Inches(3.3), Inches(5.5), Inches(3.9), SURFACE, BORDER)
txt(s, Inches(7.2), Inches(3.5), Inches(5.0), Inches(0.4), "Why AMD for AI Inference?", sz=16, color=WHITE, bold=True)
reasons = [
    "✓  ROCm ecosystem mature for LLM inference",
    "✓  vLLM native support for AMD GPUs",
    "✓  Cost-effective vs. NVIDIA for inference workloads",
    "✓  Gemma 3N 27B runs efficiently on MI300X",
    "✓  High VRAM for large context windows",
]
multi(s, Inches(7.2), Inches(4.1), Inches(5.0), Inches(2.0), reasons, sz=12, color=DIM)
rect(s, Inches(7.2), Inches(5.8), Inches(5.0), Inches(1.0), SURFACE, RGBColor(0x20, 0x20, 0x20))
multi(s, Inches(7.5), Inches(5.9), Inches(4.6), Inches(0.8), [
    "# One env var. That's it.",
    "GEMMA_API_URL=http://AMD_GPU:8000/v1/chat/completions"
], sz=10, color=DIM)

# ═══════════════════════════════════════
# SLIDE 8: TECHNICAL HIGHLIGHTS
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
eyebrow(s, "ENGINEERING")
title_block(s, "Technical Highlights")

highlights = [
    ("🛡️", "Robust JSON Extraction", "LLMs don't always return clean JSON. Our extractor strips markdown fences, then uses brace-depth tracking to find the JSON object — even in messy responses."),
    ("🔄", "Model Fallback Chain", "Primary model (Gemma 3N 27B) fails? Automatic failover to GPT-OSS 120B. Zero downtime, no user-facing errors. Graceful degradation built in."),
    ("⚡", "Parallel Signal Fetch", "All 3 forex pairs analyzed simultaneously via Promise.all. Combined with caching, dashboard loads in under 200ms on cache hit."),
    ("🧩", "Prompt Engineering", "14 technical indicators injected into a structured prompt template. The AI receives market context, not just price — producing informed analysis."),
    ("📡", "Live Rate Streaming", "Exchange rates auto-refresh every 60 seconds. The current rate is injected into each analysis prompt — AI always sees real-time data."),
    ("🎨", "Canvas Animation", "Custom particle system with 80 floating nodes, animated candlestick chart, and sine waves — all GPU-accelerated via requestAnimationFrame."),
]
for i, (icon, title, desc) in enumerate(highlights):
    col = i % 3
    row = i // 3
    left = Inches(0.9 + col * 4.0)
    top = Inches(2.8 + row * 2.2)
    rect(s, left, top, Inches(3.7), Inches(1.9), SURFACE, BORDER)
    txt(s, left + Inches(0.25), top + Inches(0.15), Inches(0.4), Inches(0.4), icon, sz=18)
    txt(s, left + Inches(0.3), top + Inches(0.55), Inches(3.2), Inches(0.35), title, sz=13, color=WHITE, bold=True)
    txt(s, left + Inches(0.3), top + Inches(0.9), Inches(3.2), Inches(0.8), desc, sz=10, color=DIM)

# ═══════════════════════════════════════
# SLIDE 9: COMPARISON
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
eyebrow(s, "DIFFERENTIATION")
title_block(s, "SignalSniper vs. Traditional Signals")

table_data = [
    ("Feature", "Traditional Signal Services", "SignalSniper AI"),
    ("Output", '"BUY EUR/USD"', "Full trade thesis with 17 structured fields"),
    ("Reasoning", "None", "AI-explained sentiment + numbered reasons"),
    ("Risk Context", 'Generic "medium"', "Dynamic SL/TP, R:R ratio, support/resistance"),
    ("Data Input", "Price only", "14 technical indicators + live rates"),
    ("AI Models", "None (rule-based)", "Dual pipeline: Fireworks + AMD GPU (Gemma)"),
    ("Explainability", "Black box", "Transparent reasoning, risk warnings, invalidation"),
    ("Latency", "Delayed alerts", "Real-time with 1hr intelligent caching"),
]

rows = len(table_data)
cols = 3
tbl_shape = s.shapes.add_table(rows, cols, Inches(1.0), Inches(3.0), Inches(11.3), Inches(4.0))
table = tbl_shape.table
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.3)
table.columns[2].width = Inches(5.0)

for i, row_data in enumerate(table_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = ""
        tf = cell.text_frame
        tf.paragraphs[0].text = cell_text

        p = tf.paragraphs[0]
        if i == 0:
            p.font.size = Pt(9)
            p.font.color.rgb = MUTED
            p.font.bold = True
            p.font.name = "Consolas"
        else:
            p.font.size = Pt(11)
            p.font.name = "Calibri"
            if j == 2:
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif j == 0:
                p.font.color.rgb = WHITE
                p.font.bold = True
            else:
                p.font.color.rgb = DIM

        cell.fill.solid()
        if i == 0:
            cell.fill.fore_color.rgb = BG
        else:
            cell.fill.fore_color.rgb = SURFACE if i % 2 == 0 else BG

# ═══════════════════════════════════════
# SLIDE 10: ROADMAP
# ═══════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
eyebrow(s, "ROADMAP")
title_block(s, "What's Next")

phases = [
    ("Phase 1 — Now", "3 major forex pairs, dual AI pipeline, live rates, structured trade cards with full AI-explained analysis.", GREEN),
    ("Phase 2 — Expand", "Add 20+ forex pairs, crypto support, multi-timeframe analysis, backtesting engine, and user-defined risk profiles.", CYAN),
    ("Phase 3 — Scale", "Real-time WebSocket price feeds, push notifications, portfolio tracking, and multi-AMD-GPU inference cluster.", PURPLE),
]
for i, (phase, desc, pc) in enumerate(phases):
    top = Inches(3.2 + i * 1.4)
    rect(s, Inches(0.9), top, Inches(5.8), Inches(1.2), SURFACE, BORDER)
    accent_bar(s, Inches(0.9), top, Inches(1.2), pc)
    txt(s, Inches(1.2), top + Inches(0.12), Inches(3), Inches(0.3), phase, sz=10, color=pc, font="Consolas")
    txt(s, Inches(1.2), top + Inches(0.45), Inches(5.2), Inches(0.6), desc, sz=12, color=DIM)

# Closing card
rect(s, Inches(7.2), Inches(3.0), Inches(5.3), Inches(4.2), SURFACE, BORDER)
txt(s, Inches(7.2), Inches(3.6), Inches(5.3), Inches(0.5), "◎  SignalSniper AI", sz=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(7.2), Inches(4.6), Inches(5.3), Inches(1.2), '"Trade smarter\nthrough understanding."', sz=20, color=DIM, font="Georgia", align=PP_ALIGN.CENTER)
rect(s, Inches(8.2), Inches(6.0), Inches(3.8), Inches(0.35), SURFACE, GREEN)
txt(s, Inches(8.4), Inches(6.02), Inches(3.6), Inches(0.3), "● github.com/public321-ai/signalsniper-ai", sz=9, color=GREEN, font="Consolas", align=PP_ALIGN.CENTER)
txt(s, Inches(8.2), Inches(6.5), Inches(3.8), Inches(0.3), "Built for AMD Hackathon 2026", sz=10, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ──
output = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SignalSniper_AI_Hackathon.pptx")
prs.save(output)
print(f"Saved: {output}")
